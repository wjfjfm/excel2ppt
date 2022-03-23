from pptx import Presentation
import pandas as pd

print('openning data.xlsx ...')
df = pd.read_excel('data.xlsx', sheet_name=0)

# parse vars from excel
var_map = {}
for loc in df:
    print('found var:', loc)
    var_map[loc] = []

if 'filename' not in var_map:
    print('Error: filename request!')
    exit(1)

print('openning template.pptx ...')
prs = Presentation('template.pptx')

print('Done')

# parse all placeholders in template ppt
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.text in var_map:
            var_map[shape.text].append(shape)

# generate pptx(s)
for index, row in df.iterrows():
    print('generating {}.pptx'.format(row['filename']))
    for (key, value) in var_map.items():
        for shape in value:
            shape.text = str(row[key]).encode('utf-8')

    prs.save('{}.pptx'.format(row['filename']))
    print('Done!')

print('Everything Done, exiting ...')
exit(0)

